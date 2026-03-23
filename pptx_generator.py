from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
import io

def generate_pptx(data):
    """
    Data: The UIOutput Pydantic object from your CrewAI task.
    """
    prs = Presentation()

    # --- Slide 1: The Reading (Title) ---
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    slide.shapes.title.text = f"The Tarot of Strategy: {data.account_name}"
    slide.placeholders[1].text = f"Incumbent: {data.competitor_name}\nDisruption Propensity: {data.displacement_score}%"

    # --- Slide 2: The Stale Incumbent (Problem Canvas) ---
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = "The Shadow: Incumbent Gaps"
    tf = slide.placeholders[1].text_frame
    tf.text = f"Specific Pain Point: {data.canvas.specific_pain_point}"
    p = tf.add_paragraph()
    p.text = f"Priority Rank: {data.priority_ranking}/5"

    # --- Slides 3-5: The Agilisium Arcana (Solutions) ---
    for s in data.slides:
        slide = prs.slides.add_slide(prs.slide_layouts[1])
        slide.shapes.title.text = s.title
        
        # Hero Metric (The Agilisium Angle)
        txBox = slide.shapes.add_textbox(Inches(1), Inches(1.5), Inches(8), Inches(1))
        p = txBox.text_frame.paragraphs[0]
        p.text = s.hero_metric
        p.font.bold = True
        p.font.size = Pt(32)
        p.alignment = PP_ALIGN.CENTER
        
        # Key Statements
        tf = slide.placeholders[1].text_frame
        for point in s.key_points:
            p = tf.add_paragraph()
            p.text = point
            p.level = 0

    # Save to binary for Streamlit download
    binary_output = io.BytesIO()
    prs.save(binary_output)
    binary_output.seek(0)
    return binary_output
